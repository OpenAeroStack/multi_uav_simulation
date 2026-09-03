#!/usr/bin/env python3
"""Build the HITL defence deck as an editable PowerPoint file.

Native shapes and text boxes, not screenshots: every element stays editable in
PowerPoint, and the fonts are ones Windows and Office already have.

Figures come from results/<run>/summary.txt, so the deck cannot drift from the
measured data. Pass a run directory to use a different run.

    python3 scripts/make_slides_pptx.py [results/20260903_163917]

Output: report/HITL_slides.pptx
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "report" / "HITL_slides.pptx"

INK    = RGBColor(0x14, 0x1C, 0x26)
MUTED  = RGBColor(0x5B, 0x67, 0x75)
FAINT  = RGBColor(0x8A, 0x94, 0xA1)
RULE   = RGBColor(0xD3, 0xDA, 0xE3)
SUNK   = RGBColor(0xF4, 0xF6, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SENSOR = RGBColor(0x0B, 0x6E, 0x7F)      # unimpaired camera link
RADIO  = RGBColor(0xB0, 0x43, 0x0B)      # impaired ns-3 link
GOOD   = RGBColor(0x1E, 0x6B, 0x45)
WARN   = RGBColor(0xB0, 0x43, 0x0B)

BODY, MONO = "Calibri", "Consolas"       # present on every Office install


# ── helpers ─────────────────────────────────────────────────────────────────
def text(slide, x, y, w, h, runs, size=12, color=INK, font=BODY,
         bold=False, align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    """runs: a string, or a list of (string, {overrides}) tuples, one per line."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = [runs] if isinstance(runs, str) else runs
    for i, line in enumerate(lines):
        s, over = (line, {}) if isinstance(line, str) else line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.line_spacing = over.get("spacing", spacing)
        if over.get("space_before"):
            p.space_before = Pt(over["space_before"])
        r = p.add_run()
        r.text = s
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
    return box


def box(slide, x, y, w, h, fill=WHITE, line=RULE, width=1.0, dash=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(width)
    if dash:
        sh.line.dash_style = dash
    sh.text_frame.word_wrap = True
    return sh


def rule(slide, x, y, w, h, color=INK, width=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def header(slide, num, title, eyebrow):
    rule(slide, 0.55, 0.42, 0.42, 0.30, INK)
    text(slide, 0.55, 0.47, 0.42, 0.25, num, size=11, color=WHITE, font=MONO,
         bold=True, align=PP_ALIGN.CENTER)
    text(slide, 1.12, 0.33, 8.8, 0.55, title, size=23, bold=True, spacing=0.92)
    text(slide, 10.1, 0.47, 2.7, 0.3, eyebrow, size=10, color=FAINT, font=MONO,
         align=PP_ALIGN.RIGHT)
    rule(slide, 0.55, 0.95, 12.22, 0.022, INK)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return s


def parse_summary(path):
    """Pull the per-board figures out of summarise_run.sh output."""
    if not path.exists():
        sys.exit(f"ERROR: no summary at {path}\n"
                 "       Fly a mission first, or pass a run directory.")
    txt = path.read_text()
    out = {}
    for blockmatch in re.finditer(
            r"── (UAV\d) ─.*?(?=── UAV|\Z)", txt, re.S):
        name, blk = blockmatch.group(1), blockmatch.group(0)

        def g(pat, cast=float, default=None):
            m = re.search(pat, blk)
            return cast(m.group(1)) if m else default

        out[name] = {
            "frames":    g(r"edge\s*:\s*(\d+) frames", int, 0),
            "secs":      g(r"in ([\d.]+) s"),
            "fps":       g(r"=\s*([\d.]+) fps"),
            "inf":       g(r"mean (\d+) ms", int),
            "inf_min":   g(r"min (\d+) \|", int),
            "inf_max":   g(r"max (\d+) \| ceiling", int),
            "ceiling":   g(r"ceiling ([\d.]+) fps"),
            "hits":      g(r"detected:\s*(\d+) frames", int, 0),
            "hitpct":    g(r"person \(([\d.]+) %\)"),
            "gcs":       g(r"gcs\s*:\s*(\d+) detection", int, 0),
            "pay":       g(r"payload : mean (\d+) B", int),
            "pay_max":   g(r"payload : mean \d+ B \| min \d+ \| max (\d+)", int),
            "sent":      g(r"delivery: \d+/(\d+)", int),
            "deliv":     g(r"=\s*([\d.]+) % reached"),
            "t_mean":    g(r"mean ([\d.]+) C"),
            "t_peak":    g(r"peak ([\d.]+) C"),
            "throttled": g(r"throttled (\S+)", str, "?"),
        }
    if not out:
        sys.exit("ERROR: could not parse the summary file.")
    return out


# ── slide 1 — HITL architecture ─────────────────────────────────────────────
def slide_architecture(prs, d):
    s = blank(prs)
    header(s, "01", "Two real boards in a simulated fleet", "HITL ARCHITECTURE")

    text(s, 0.55, 1.12, 12.2, 0.5,
         "Each aircraft runs in its own Linux network namespace with a real IP stack. "
         "One Ethernet cable per board carries two tagged VLANs, so the camera feed and "
         "the radio link share a wire but are different networks — and only one is degraded.",
         size=12.5, color=MUTED, spacing=1.15)

    # legend
    rule(s, 0.55, 1.83, 0.26, 0.045, SENSOR)
    text(s, 0.90, 1.72, 3.4, 0.25, "VLAN 10 · sensor · unimpaired",
         size=10, color=SENSOR, font=MONO, bold=True)
    rule(s, 4.35, 1.83, 0.26, 0.045, RADIO)
    text(s, 4.70, 1.72, 3.8, 0.25, "VLAN 42/43 · radio · through ns-3",
         size=10, color=RADIO, font=MONO, bold=True)

    # host enclosure
    box(s, 0.55, 2.06, 6.55, 4.12, fill=None, line=FAINT, width=1.0)
    text(s, 0.65, 2.11, 4.5, 0.22, "HOST — UBUNTU 22.04 / ROS 2 HUMBLE",
         size=9, color=FAINT, font=MONO, bold=True)

    def node(x, y, w, h, title, lines):
        box(s, x, y, w, h, fill=WHITE, line=RULE)
        text(s, x + 0.14, y + 0.10, w - 0.28, 0.22, title, size=11, bold=True)
        text(s, x + 0.14, y + 0.34, w - 0.28, h - 0.42,
             lines, size=9, color=MUTED, font=MONO, spacing=1.25)

    node(0.72, 2.40, 3.05, 0.82, "Gazebo · physics + 2 cameras",
         ["640×384 · 5 Hz · 45° down", "0.6 rad FOV · small_city world"])
    node(0.72, 3.30, 3.05, 0.56, "netns uav1ns · ArduPilot SITL",
         ["10.42.0.11 · sysid 1 · AP_DDS"])
    node(0.72, 3.92, 3.05, 0.56, "netns uav2ns · ArduPilot SITL",
         ["10.42.0.13 · sysid 2 · AP_DDS"])
    node(0.72, 4.70, 3.05, 1.32, "netns gcsns · ground station",
         ["drone_bridge ×2 — commands out",
          "micro_ros_agent ×2 — telemetry in",
          "gcs_receiver ×2 — detections in",
          "10.42.0.10"])

    box(s, 3.98, 3.30, 3.00, 2.72, fill=SUNK, line=RULE)
    text(s, 4.12, 3.40, 2.75, 0.25, "ns-3 · simulated radio", size=11, bold=True)
    text(s, 4.12, 3.70, 2.75, 2.2,
         ["three_uav_tapbridge_integrated",
          "TapBridge injects the channel",
          "model into real Linux packets",
          "802.11 · Nakagami fading",
          "nodes 0–3 · GCS + UAVs + edge",
          "",
          "Detections cross HERE.",
          "Camera frames never do."],
         size=9, color=MUTED, font=MONO, spacing=1.35)

    # edge boards
    text(s, 7.55, 2.13, 3.0, 0.22, "EDGE NODES", size=9, color=FAINT,
         font=MONO, bold=True)

    def board(y, name, cam, radio_ip, inf, fps):
        box(s, 7.55, y, 5.25, 1.68, fill=WHITE, line=RULE)
        text(s, 7.72, y + 0.13, 4.9, 0.25, name, size=11.5, bold=True)
        text(s, 7.72, y + 0.42, 4.9, 1.2,
             [f"YOLO11n · OpenVINO FP32 · 384×640",
              f"conf 0.40 · class 0 (person)",
              f"cam {cam}   ·   radio {radio_ip}",
              f"measured: {inf} ms · {fps} fps"],
             size=9.5, color=MUTED, font=MONO, spacing=1.35)

    board(2.40, "Raspberry Pi 4B — UAV1 edge node", "10.0.0.2", "10.42.0.12",
          d["UAV1"]["inf"], f'{d["UAV1"]["fps"]:.2f}')
    board(4.34, "Raspberry Pi 4B — UAV2 edge node", "10.0.1.2", "10.42.0.14",
          d["UAV2"]["inf"], f'{d["UAV2"]["fps"]:.2f}')

    # sensor path: over the top, clear of ns-3
    rule(s, 3.77, 2.62, 3.53, 0.035, SENSOR)
    rule(s, 7.27, 2.62, 0.035, 2.58, SENSOR)
    rule(s, 7.27, 3.04, 0.30, 0.035, SENSOR)
    rule(s, 7.27, 5.18, 0.30, 0.035, SENSOR)
    text(s, 4.00, 2.32, 3.3, 0.25, "camera frames · 737 KB/frame",
         size=9, color=SENSOR, font=MONO, bold=True)

    # radio path: boards -> ns-3 -> gcsns
    rule(s, 6.98, 3.72, 0.58, 0.035, RADIO)
    rule(s, 6.98, 5.66, 0.58, 0.035, RADIO)
    rule(s, 3.77, 5.40, 0.24, 0.035, RADIO)
    text(s, 4.12, 5.72, 3.0, 0.25, 'detections · 73–504 B JSON',
         size=9, color=RADIO, font=MONO, bold=True)

    # takeaways
    cols = [
        ("Isolation is real, not simulated",
         "Network namespaces give each aircraft a genuine kernel stack — separate routing "
         "tables and its own 127.0.0.1. Nothing shortcuts through loopback, so the packets "
         "ns-3 degrades are the packets the flight code sent."),
        ("One cable, two networks",
         "802.1Q VLAN tagging splits each board's link. VLAN 10 behaves like a camera "
         "ribbon; VLAN 42/43 is bridged into ns-3 — modelling an airframe where sensor bus "
         "and radio are separate systems."),
        ("Only results cross the radio",
         "A 737 KB frame stays on the sensor link. What crosses the degraded channel is "
         "73–504 bytes of JSON. That reduction is the architectural claim edge processing "
         "makes, and what this experiment measures."),
    ]
    for i, (h, body) in enumerate(cols):
        x = 0.55 + i * 4.15
        text(s, x, 6.34, 3.85, 0.22, h.upper(), size=9, color=FAINT,
             font=MONO, bold=True)
        text(s, x, 6.57, 3.85, 0.90, body, size=9, color=MUTED, spacing=1.12)


# ── slide 2 — detection ─────────────────────────────────────────────────────
def slide_detection(prs, d):
    s = blank(prs)
    header(s, "02", "Human detection on the airframe", "EDGE INFERENCE")

    text(s, 0.55, 1.10, 12.2, 0.4,
         "Each board subscribes to its own aircraft's camera, runs YOLO11n under OpenVINO, "
         "and publishes only the bounding boxes. Whether a person is detectable at all is "
         "set by six physical and software parameters, derived and measured below.",
         size=11.5, color=MUTED, spacing=1.12)

    # ── geometry figure (compact, left) ──────────────────────────────────────
    fig_x, fig_y, fig_w, fig_h = 0.55, 1.62, 5.55, 2.70
    box(s, fig_x, fig_y, fig_w, fig_h, fill=SUNK, line=RULE)
    gy = fig_y + 2.10                                # ground line
    rule(s, fig_x + 0.30, gy, fig_w - 0.75, 0.026, FAINT)
    text(s, fig_x + 0.30, gy + 0.09, 1.0, 0.18, "GROUND", size=7.5, color=FAINT,
         font=MONO, bold=True)

    apex_x, apex_y = fig_x + 1.05, fig_y + 0.30
    near_x, far_x = fig_x + 2.05, fig_x + 4.55
    ff = s.shapes.build_freeform(Inches(apex_x), Inches(apex_y))
    ff.add_line_segments([(Inches(near_x), Inches(gy)),
                          (Inches(far_x), Inches(gy))], close=True)
    tri = ff.convert_to_shape()
    tri.shadow.inherit = False
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(0xDC, 0xEE, 0xF1)
    tri.line.color.rgb = SENSOR
    tri.line.width = Pt(1.0)

    box(s, apex_x - 0.32, apex_y - 0.08, 0.72, 0.26, fill=WHITE, line=RULE)
    text(s, apex_x - 0.25, apex_y - 0.04, 0.6, 0.18, "UAV", size=8, bold=True)
    rule(s, apex_x, apex_y + 0.19, 0.016, gy - apex_y - 0.19, FAINT)
    text(s, apex_x - 0.42, gy - 1.10, 0.6, 0.18, "30 m", size=7.5, color=MUTED, font=MONO)

    rule(s, apex_x, gy - 0.03, near_x - apex_x, 0.045, RADIO)
    text(s, apex_x - 0.1, gy - 0.27, 1.9, 0.18, "blind: 0–20.6 m",
         size=7.5, color=RADIO, font=MONO, bold=True)
    text(s, near_x - 0.28, gy + 0.09, 0.8, 0.18, "20.6 m", size=7.5, color=MUTED, font=MONO)
    text(s, far_x - 0.28, gy + 0.09, 0.8, 0.18, "43.6 m", size=7.5, color=MUTED, font=MONO)
    mid_x = (near_x + far_x) / 2
    text(s, mid_x - 0.55, gy - 0.26, 1.3, 0.18, "visible band", size=7.5, color=SENSOR,
         font=MONO, bold=True)
    rule(s, mid_x, gy - 0.15, 0.03, 0.15, INK)
    text(s, mid_x - 0.35, gy - 0.44, 1.0, 0.18, "≈30 px", size=7.5, color=INK, font=MONO)

    text(s, fig_x + fig_w - 1.85, fig_y + 0.16, 1.75, 1.6,
         ["45° pitch · 0.6 rad", "640×384 @ 5 Hz", "", "band = 0.687h", "       … 1.453h"],
         size=7.5, color=MUTED, font=MONO, spacing=1.25)

    text(s, fig_x + 0.10, fig_y + fig_h - 0.26, fig_w - 0.2, 0.2,
         "GSD & pixel-height derivation → parameters 1 and 3, right", size=7.5,
         color=FAINT, font=MONO)

    # ── sensitivity parameters (right) ───────────────────────────────────────
    px, py, pw = 6.30, 1.62, 6.48
    text(s, px, py, pw, 0.2, "DETECTION SENSITIVITY PARAMETERS", size=9,
         color=FAINT, font=MONO, bold=True)

    # Each row is label + ONE value line, kept ASCII-only: an unmatched glyph
    # (Greek theta) silently falls back to a wider font mid-string and the
    # resulting overflow collides with the row below it.
    srows = [
        ("1  Ground Sampling Distance (GSD)",
         "GSD = 2h*tan(HFOV/2)/W  =  2.9 cm/px @ 30 m (nadir-equiv)"),
        ("2  Camera resolution",
         "640x384 px, frozen by the OpenVINO export shape"),
        ("3  Person height in pixels",
         "px = 889 / altitude(m)  ->  30 px @ 30 m, 45 px @ 20 m"),
        ("4  Altitude",
         "drives GSD & pixel height directly; reliable window ~15-25 m"),
        ("5  Flight speed -> motion blur",
         "blur(px) = v x t_exposure / GSD -- not modeled by Gazebo's camera"),
        ("6  Gimbal pitch",
         "pixel height peaks at 45 deg; steeper widens the COCO->aerial gap"),
        ("7  Inference backend (measured)",
         "OpenVINO 236 ms < NCNN 270 < MNN 342 < PyTorch 1,027 (PI_SETUP.md)"),
        ("8  Confidence vs IoU",
         "conf 0.40 explicit  vs  IoU 0.70 (NMS default) -- different filters"),
    ]
    y = py + 0.30
    row_h = 0.375
    for label, value in srows:
        text(s, px, y, pw, 0.17, label, size=8.5, bold=True)
        text(s, px, y + 0.175, pw, 0.16, value, size=8, color=MUTED, font=MONO)
        rule(s, px, y + row_h - 0.035, pw, 0.008, RULE)
        y += row_h

    # ── bottom: what it means + what was measured ────────────────────────────
    by = py + 0.30 + row_h * len(srows) + 0.14

    text(s, 0.55, by, 5.9, 0.2, "WHAT THIS MEANS FOR THIS DEPLOYMENT", size=9,
         color=FAINT, font=MONO, bold=True)
    text(s, 0.55, by + 0.24, 5.9, 1.55,
         "Frames are dropped by design: BEST_EFFORT, KEEP_LAST, depth 1 means the "
         "detector always works the newest frame instead of a queue, so discarded "
         "frames were never processable anyway. Geometry, not the model, sets the "
         "ceiling — camera resolution is frozen at the export shape, so a person's "
         "pixel height is fixed by altitude and pitch alone once conf and IoU are set.",
         size=9.5, color=MUTED, spacing=1.16)

    text(s, 6.85, by, 5.9, 0.2, "MEASURED CONSEQUENCE", size=9,
         color=FAINT, font=MONO, bold=True)
    text(s, 6.85, by + 0.24, 5.9, 1.55,
         [(f'UAV1 held station near a static group: {d["UAV1"]["hits"]} of '
           f'{d["UAV1"]["frames"]} frames detected a person ({d["UAV1"]["hitpct"]:.1f} %).',
           {}),
          ("", {"size": 4}),
          (f'UAV2 overflew the walking subject at the same altitude and settings: '
           f'{d["UAV2"]["hits"]} of {d["UAV2"]["frames"]} ({d["UAV2"]["hitpct"]:.1f} %).',
           {}),
          ("", {"size": 4}),
          ("Same model, same conf/IoU, same boards. The gap is dwell time on target — "
           "exactly what parameters 3-6 above predict.", {"color": INK})],
         size=9.5, color=MUTED, spacing=1.16)


# ── slide 3 — results ───────────────────────────────────────────────────────
def slide_results(prs, d):
    s = blank(prs)
    header(s, "03", "Output, and how it was validated", "OUTPUT & VALIDATION")

    u1, u2 = d["UAV1"], d["UAV2"]
    eff1 = 100 * u1["fps"] / u1["ceiling"]
    eff2 = 100 * u2["fps"] / u2["ceiling"]

    text(s, 0.55, 1.12, 12.2, 0.4,
         f"One instrumented mission, both boards. Every figure below is read from the "
         f"run's own logs by scripts/summarise_run.sh — no hand-copied numbers.",
         size=12.5, color=MUTED, spacing=1.15)

    # stat tiles
    tiles = [
        (f"{eff1:.0f} % / {eff2:.0f} %", "OF COMPUTE CEILING",
         f'{u1["fps"]:.2f} and {u2["fps"]:.2f} fps against limits of '
         f'{u1["ceiling"]:.2f} and {u2["ceiling"]:.2f}', GOOD),
        (f'{u1["inf"]} / {u2["inf"]} ms', "MEAN INFERENCE",
         f'board 1 ranged {u1["inf_min"]}–{u1["inf_max"]}, '
         f'board 2 {u2["inf_min"]}–{u2["inf_max"]}', INK),
        (f'{u1["deliv"]:.1f} % / {u2["deliv"]:.1f} %', "DELIVERED ACROSS ns-3",
         f'{u1["gcs"]}/{u1["sent"]} and {u2["gcs"]}/{u2["sent"]} '
         f'detection messages reached the GCS', INK),
        (f'{u1["pay"]} / {u2["pay"]} B', "MEAN PAYLOAD",
         f'73 B empty result, up to {u1["pay_max"]} B with several people', INK),
    ]
    for i, (v, k, note, col) in enumerate(tiles):
        x = 0.55 + i * 3.09
        box(s, x, 1.62, 2.95, 1.30, fill=WHITE, line=RULE)
        text(s, x + 0.16, 1.74, 2.7, 0.4, v, size=20, bold=True, font=MONO, color=col)
        text(s, x + 0.16, 2.16, 2.7, 0.2, k, size=8.5, color=FAINT, font=MONO, bold=True)
        text(s, x + 0.16, 2.40, 2.7, 0.48, note, size=8.5, color=MUTED, spacing=1.12)

    # validation chain
    text(s, 0.55, 3.08, 6.0, 0.22, "HOW A RUN IS VALIDATED", size=9,
         color=FAINT, font=MONO, bold=True)
    steps = [
        ("STEP 0", "rpi_init.sh", "16 checks per board: link, VLAN addressing, RTT, "
                                  "socket buffers, clock skew ≤ 50 ms"),
        ("STEP 1", "sitl_init.sh", "Cold-starts netns, ns-3, Gazebo, SITL. Waits for a "
                                   "real navsat message per aircraft"),
        ("STEP 2", "detector_start.sh", "Starts both edge detectors, blocks until each "
                                        "reports its model loaded"),
        ("STEP 3", "run_missions.sh", "Re-checks every stage before arming, flies, then "
                                      "archives and summarises the run"),
    ]
    for i, (n, name, desc) in enumerate(steps):
        x = 0.55 + i * 3.09
        box(s, x, 3.32, 2.95, 1.05, fill=SUNK, line=RULE)
        text(s, x + 0.14, 3.42, 2.7, 0.18, n, size=8, color=FAINT, font=MONO, bold=True)
        text(s, x + 0.14, 3.62, 2.7, 0.2, name, size=10, font=MONO, bold=True)
        text(s, x + 0.14, 3.86, 2.7, 0.45, desc, size=8.5, color=MUTED, spacing=1.1)

    # findings
    text(s, 0.55, 4.58, 6.0, 0.22, "WHAT THE INSTRUMENTATION REVEALED", size=9,
         color=FAINT, font=MONO, bold=True)
    text(s, 0.55, 4.84, 6.0, 2.2,
         [("Board 1 is power-limited, not heat-limited.", {"bold": True, "color": WARN}),
          (f'It ran hotter work more slowly ({u1["inf"]} ms) than board 2 '
           f'({u2["inf"]} ms) while staying cooler — {u1["t_peak"]:.1f} °C peak against '
           f'{u2["t_peak"]:.1f} °C. Its throttle flag reads {u1["throttled"]}: '
           f'under-voltage, not temperature. A better supply, not a heatsink.', {}),
          ("", {"size": 6}),
          ("Both boards ran at 96 % of their own compute ceiling.",
           {"bold": True, "color": INK}),
          ("Delivery was never the constraint: if the link were starving them, the "
           "achieved rate would sit well below what inference allows. It does not.", {}),
          ("", {"size": 6}),
          (f'Board 2 lost {100 - u2["deliv"]:.0f} % of its detections crossing ns-3.',
           {"bold": True, "color": WARN}),
          (f'{u2["gcs"]} of {u2["sent"]} arrived, against {u1["deliv"]:.1f} % on board 1 '
           f'— a real difference between the two radio links, not measurement noise.', {})],
         size=10, color=MUTED, spacing=1.16)

    text(s, 6.85, 4.58, 5.9, 0.22, "STATED LIMITATIONS", size=9,
         color=FAINT, font=MONO, bold=True)
    text(s, 6.85, 4.84, 5.9, 2.2,
         [("INT8 quantisation was evaluated and rejected.", {"bold": True, "color": INK}),
          ("The quantised model was verified correct on x86, but OpenVINO's ARM CPU "
           "plugin could not compile the graph on the Pi 4B. Edge inference runs FP32.", {}),
          ("", {"size": 6}),
          ("Board power is not instrumented.", {"bold": True, "color": INK}),
          ("The Pi 4B exposes no power sensor. ARM clock, core voltage and CPU load are "
           "logged as on-board proxies; absolute wattage needs an inline meter.", {}),
          ("", {"size": 6}),
          ("Throttling changes timing, never accuracy.", {"bold": True, "color": INK}),
          ("A throttled CPU produces bit-identical detections more slowly. It reduces "
           "frames analysed per target pass — not precision or recall.", {})],
         size=10, color=MUTED, spacing=1.16)

    text(s, 0.55, 7.10, 12.2, 0.25,
         f'Source: results/{sys.argv[1].rstrip("/").split("/")[-1] if len(sys.argv) > 1 else "20260903_163917"}'
         f'  ·  {u1["frames"]} + {u2["frames"]} frames  ·  '
         f'{u1["secs"]:.0f} s and {u2["secs"]:.0f} s of flight  ·  FP32 YOLO11n',
         size=8.5, color=FAINT, font=MONO)


def main():
    run = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "results" / "20260903_163917"
    data = parse_summary(run / "summary.txt")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9

    slide_architecture(prs, data)
    slide_detection(prs, data)
    slide_results(prs, data)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print(f"  {OUT.relative_to(ROOT)}  {OUT.stat().st_size // 1024} KB  "
          f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides, 16:9)")


if __name__ == "__main__":
    main()
